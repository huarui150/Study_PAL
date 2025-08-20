"""
学习资料管理服务 - 修复版
负责文件上传、处理、存储和文本提取
"""
import os
import hashlib
import pymysql
from typing import Dict, List, Any, Tuple, Optional
from werkzeug.datastructures import FileStorage
from werkzeug.utils import secure_filename
import json
from auth import get_db_connection


class MaterialService:
    """学习资料服务类"""
    
    UPLOAD_FOLDER = 'uploads'
    ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx', 'doc', 'docx', 'png', 'jpg', 'jpeg', 'gif', 'txt'}
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    
    def __init__(self):
        # 确保上传目录存在
        os.makedirs(self.UPLOAD_FOLDER, exist_ok=True)
    
    @staticmethod
    def allowed_file(filename: str) -> bool:
        """检查文件类型是否允许"""
        return '.' in filename and \
               filename.rsplit('.', 1)[1].lower() in MaterialService.ALLOWED_EXTENSIONS
    
    @staticmethod
    def get_file_type(filename: str) -> str:
        """根据文件名获取文件类型"""
        if not filename or '.' not in filename:
            return 'unknown'
            
        ext = filename.rsplit('.', 1)[1].lower()
        if ext == 'pdf':
            return 'pdf'
        elif ext in ['ppt', 'pptx']:
            return 'ppt'
        elif ext in ['doc', 'docx']:
            return 'doc'
        elif ext in ['png', 'jpg', 'jpeg', 'gif']:
            return 'image'
        elif ext == 'txt':
            return 'txt'
        return 'unknown'
    
    @staticmethod
    def calculate_file_hash(file_content: bytes) -> str:
        """计算文件哈希值"""
        return hashlib.sha256(file_content).hexdigest()
    
    @staticmethod
    def extract_text_from_file(file_path: str, file_type: str) -> str:
        """从文件中提取文本内容"""
        try:
            print(f"[DEBUG] 开始提取文本: {file_path}, 类型: {file_type}")
            
            if file_type == 'txt':
                return MaterialService._extract_from_txt(file_path)
            elif file_type == 'pdf':
                return MaterialService._extract_from_pdf(file_path)
            elif file_type == 'ppt':
                return MaterialService._extract_from_ppt(file_path)
            elif file_type == 'doc':
                return MaterialService._extract_from_doc(file_path)
            elif file_type == 'image':
                return MaterialService._extract_from_image(file_path)
            else:
                return "不支持的文件类型"
        except Exception as e:
            print(f"[ERROR] 文本提取失败: {e}")
            return f"文本提取失败: {str(e)}"
    
    @staticmethod
    def _extract_from_txt(file_path: str) -> str:
        """从TXT文件提取文本"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    return f.read()
            except:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
        except Exception as e:
            return f"TXT处理错误: {str(e)}"
    
    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        """从PDF提取文本（多引擎容错：PyMuPDF → pdfminer.six → PyPDF2）"""
        # 1) 优先使用 PyMuPDF（pymupdf）
        try:
            import fitz  # PyMuPDF
            text = ""
            with fitz.open(file_path) as doc:
                for page in doc:
                    page_text = page.get_text("text")
                    if page_text:
                        text += page_text + "\n"
            if text.strip():
                return text.strip()
        except ImportError:
            pass
        except Exception as e:
            print(f"[WARNING] PyMuPDF 提取失败: {e}")

        # 2) 尝试 pdfminer.six
        try:
            from pdfminer.high_level import extract_text as pdfminer_extract_text
            text = pdfminer_extract_text(file_path) or ""
            if text.strip():
                return text.strip()
        except ImportError:
            pass
        except Exception as e:
            print(f"[WARNING] pdfminer 提取失败: {e}")

        # 3) 回退到 PyPDF2
        try:
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file, strict=False)
                if getattr(reader, 'is_encrypted', False):
                    try:
                        reader.decrypt('')
                    except Exception:
                        pass
                for page in reader.pages:
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    except Exception:
                        continue
            if text.strip():
                return text.strip()
        except ImportError:
            return "需要安装PyMuPDF或pdfminer.six或PyPDF2库来处理PDF文件：pip install pymupdf pdfminer.six PyPDF2"
        except Exception as e:
            print(f"[WARNING] PyPDF2 提取失败: {e}")

        return "PDF文本提取失败。建议：1) 重新导出为标准PDF；2) 转换为PDF/A；或 3) 先转为TXT/图片再上传。"
    
    @staticmethod
    def _extract_from_ppt(file_path: str) -> str:
        """从PPT提取文本"""
        try:
            from pptx import Presentation
            text = ""
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n=== 幻灯片 {slide_num} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text.strip() or "PPT文件为空或无法提取文本"
        except ImportError:
            return "需要安装python-pptx库来处理PPT文件：pip install python-pptx"
        except Exception as e:
            return f"PPT处理错误: {str(e)}"
    
    @staticmethod
    def _extract_from_doc(file_path: str) -> str:
        """从Word文档提取文本"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            return text.strip() or "Word文件为空或无法提取文本"
        except ImportError:
            return "需要安装python-docx库来处理Word文件：pip install python-docx"
        except Exception as e:
            return f"Word处理错误: {str(e)}"
    
    @staticmethod
    def _extract_from_image(file_path: str) -> str:
        """从图片提取文本（OCR）"""
        return "图片已上传，OCR功能开发中，暂时无法提取文本内容"
    
    @staticmethod
    def upload_material(user_id: int, course_id: int, file: FileStorage, title: str = None) -> Tuple[Dict[str, Any], int]:
        """上传学习资料"""
        try:
            print(f"[DEBUG] 开始上传文件 - user_id: {user_id}, course_id: {course_id}")
            print(f"[DEBUG] 文件信息: {file.filename if file else 'None'}")
            
            # 1. 验证输入参数
            if not user_id or not course_id:
                return {"error": "用户ID和课程ID不能为空"}, 400
                
            if not file or not file.filename:
                return {"error": "没有选择文件"}, 400
            
            if not MaterialService.allowed_file(file.filename):
                return {"error": f"不支持的文件类型，支持: {', '.join(MaterialService.ALLOWED_EXTENSIONS)}"}, 400
            
            # 2. 读取文件内容
            file_content = file.read()
            file.seek(0)  # 重置文件指针
            
            if len(file_content) == 0:
                return {"error": "文件为空"}, 400
                
            if len(file_content) > MaterialService.MAX_FILE_SIZE:
                return {"error": f"文件大小超过限制({MaterialService.MAX_FILE_SIZE // 1024 // 1024}MB)"}, 400
            
            # 3. 计算文件哈希和获取文件信息
            file_hash = MaterialService.calculate_file_hash(file_content)
            original_filename = file.filename or "uploaded_file"
            # 使用哈希+原扩展名作为存储文件名，支持中文标题不受影响
            _, ext = os.path.splitext(original_filename)
            ext = (ext or '').lower()
            storage_filename = f"{file_hash}{ext if ext else ''}"
            # 规范化type到数据库允许的枚举(pdf|ppt|doc|image)
            detected_type = MaterialService.get_file_type(original_filename)
            type_mapping = {
                'pdf': 'pdf',
                'ppt': 'ppt',
                'pptx': 'ppt',
                'doc': 'doc',
                'docx': 'doc',
                'txt': 'doc',
                'png': 'image',
                'jpg': 'image',
                'jpeg': 'image',
                'gif': 'image',
                'image': 'image'
            }
            file_type = type_mapping.get(detected_type, 'doc')
            material_title = title or original_filename
            
            print(f"[DEBUG] 文件处理 - 原名: {original_filename}, 存储名: {storage_filename}, 类型: {file_type}, 哈希: {file_hash[:10]}...")
            
            # 4. 检查数据库连接和课程权限
            try:
                connection = get_db_connection()
                print("[DEBUG] 数据库连接成功")
                
                with connection.cursor(pymysql.cursors.DictCursor) as cursor:
                    # 检查课程是否存在且用户有权限
                    cursor.execute(
                        "SELECT id, name FROM courses WHERE id = %s AND user_id = %s",
                        (course_id, user_id)
                    )
                    course = cursor.fetchone()
                    if not course:
                        connection.close()
                        return {"error": "课程不存在或无权限访问"}, 404
                    
                    print(f"[DEBUG] 课程验证成功: {course['name']}")
                    
                    # 检查是否已上传相同文件
                    cursor.execute(
                        "SELECT id, title FROM materials WHERE course_id = %s AND file_hash = %s",
                        (course_id, file_hash)
                    )
                    existing = cursor.fetchone()
                    if existing:
                        connection.close()
                        return {"error": f"相同文件已存在: {existing['title']}"}, 409
                    
                    print("[DEBUG] 文件唯一性检查通过")
                    
            except Exception as db_error:
                print(f"[ERROR] 数据库操作失败: {db_error}")
                return {"error": f"数据库错误: {str(db_error)}"}, 500
            
            # 5. 保存文件到磁盘
            try:
                file_path = os.path.join(MaterialService.UPLOAD_FOLDER, storage_filename)
                with open(file_path, 'wb') as f:
                    f.write(file_content)
                print(f"[DEBUG] 文件保存成功: {file_path}")
            except Exception as save_error:
                connection.close()
                return {"error": f"文件保存失败: {str(save_error)}"}, 500
            
            # 6. 提取文本内容（不阻塞主流程）
            text_content = ""
            try:
                text_content = MaterialService.extract_text_from_file(file_path, file_type)
                print(f"[DEBUG] 文本提取完成，长度: {len(text_content)}")
            except Exception as extract_error:
                print(f"[WARNING] 文本提取失败: {extract_error}")
                text_content = "文本提取失败，但文件已成功上传"
            
            # 7. 保存到数据库
            try:
                with connection.cursor() as cursor:
                    cursor.execute("""
                        INSERT INTO materials (course_id, title, type, url, file_hash)
                        VALUES (%s, %s, %s, %s, %s)
                    """, (course_id, material_title, file_type, file_path, file_hash))
                    
                    material_id = cursor.lastrowid
                    connection.commit()
                    connection.close()
                    
                    print(f"[DEBUG] 数据库记录创建成功，ID: {material_id}")
                    
            except Exception as db_save_error:
                connection.close()
                # 删除已保存的文件
                try:
                    if os.path.exists(file_path):
                        os.remove(file_path)
                except:
                    pass
                return {"error": f"数据库保存失败: {str(db_save_error)}"}, 500
            
            # 8. 返回成功结果
            return {
                "material_id": material_id,
                "title": material_title,
                "type": file_type,
                "file_hash": file_hash,
                "text_length": len(text_content),
                "text_preview": text_content[:200] + "..." if len(text_content) > 200 else text_content,
                "message": "文件上传成功"
            }, 201
            
        except Exception as e:
            print(f"[ERROR] 上传过程中出现未预期错误: {e}")
            import traceback
            traceback.print_exc()
            return {"error": f"上传失败: {str(e)}"}, 500
    
    @staticmethod
    def get_course_materials(user_id: int, course_id: int) -> Tuple[Dict[str, Any], int]:
        """获取课程的所有资料"""
        try:
            print(f"[DEBUG] 获取课程资料 - user_id: {user_id}, course_id: {course_id}")
            
            connection = get_db_connection()
            with connection.cursor(pymysql.cursors.DictCursor) as cursor:
                # 验证课程归属
                cursor.execute(
                    "SELECT name FROM courses WHERE id = %s AND user_id = %s",
                    (course_id, user_id)
                )
                course = cursor.fetchone()
                if not course:
                    connection.close()
                    return {"error": "课程不存在或无权限"}, 404
                
                # 获取资料列表（使用子查询统计卡片数，限制返回数量）
                cursor.execute("""
                    SELECT 
                        m.id, m.course_id, m.title, m.type, m.url, m.file_hash, m.uploaded_at,
                        (SELECT COUNT(*) FROM cards c WHERE c.material_id = m.id) AS card_count
                    FROM materials m
                    WHERE m.course_id = %s
                    ORDER BY m.uploaded_at DESC
                    LIMIT 500
                """, (course_id,))
                
                materials = cursor.fetchall()
                connection.close()
                
                # 格式化时间
                for material in materials:
                    if material['uploaded_at']:
                        material['uploaded_at'] = material['uploaded_at'].isoformat()
                
                print(f"[DEBUG] 找到 {len(materials)} 个资料")
                
                return {
                    "course_name": course['name'],
                    "materials": materials,
                    "total": len(materials)
                }, 200
                
        except Exception as e:
            print(f"[ERROR] 获取资料列表失败: {e}")
            return {"error": f"获取资料列表失败: {str(e)}"}, 500
    
    @staticmethod
    def delete_material(user_id: int, material_id: int) -> Tuple[Dict[str, Any], int]:
        """删除学习资料"""
        try:
            connection = get_db_connection()
            with connection.cursor() as cursor:
                # 验证资料归属
                cursor.execute("""
                    SELECT m.url, m.title FROM materials m
                    JOIN courses c ON m.course_id = c.id
                    WHERE m.id = %s AND c.user_id = %s
                """, (material_id, user_id))
                
                result = cursor.fetchone()
                if not result:
                    connection.close()
                    return {"error": "资料不存在或无权限"}, 404
                
                file_path, title = result
                
                # 删除数据库记录
                cursor.execute("DELETE FROM materials WHERE id = %s", (material_id,))
                connection.commit()
                connection.close()
                
                # 删除文件
                try:
                    if os.path.exists(file_path):
                        os.remove(file_path)
                        print(f"[DEBUG] 文件删除成功: {file_path}")
                except Exception as file_error:
                    print(f"[WARNING] 文件删除失败: {file_error}")
                
                return {"message": f"资料 '{title}' 删除成功"}, 200
                
        except Exception as e:
            print(f"[ERROR] 删除资料失败: {e}")
            return {"error": f"删除资料失败: {str(e)}"}, 500


# 对外接口函数
def handle_upload_material(user_id: int, course_id: int, file: FileStorage, title: str = None):
    """处理上传资料请求"""
    return MaterialService.upload_material(user_id, course_id, file, title)

def handle_get_materials(user_id: int, course_id: int):
    """处理获取资料列表请求"""
    return MaterialService.get_course_materials(user_id, course_id)

def handle_delete_material(user_id: int, material_id: int):
    """处理删除资料请求"""
    return MaterialService.delete_material(user_id, material_id)